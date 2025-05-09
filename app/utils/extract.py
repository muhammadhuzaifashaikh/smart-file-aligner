import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pdfplumber

def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == ".pdf":
            try:
                with pdfplumber.open(filepath) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except Exception:
                # fallback to PyMuPDF
                doc = fitz.open(filepath)
                return "\n".join(page.get_text() for page in doc)

        elif ext == ".docx":
            doc_file = docx.Document(filepath)
            return "\n".join(p.text for p in doc_file.paragraphs)

        elif ext == ".pptx":
            prs = Presentation(filepath)
            return "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""

    return ""
